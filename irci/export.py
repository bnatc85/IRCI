# irci/export.py
"""
Enhanced Export Capabilities for IRCI

Provides:
1. Excel Export - Multi-sheet workbook with all data and charts
2. CSV Batch Export - Multiple CSV files in a ZIP archive
3. PowerPoint Export - Executive summary presentation
4. PDF Export - Formatted report (via existing report_generator)
"""
from __future__ import annotations
import io
import zipfile
from datetime import datetime
from typing import Dict, List, Optional, Any
import pandas as pd
import numpy as np

from .logging import get_logger

log = get_logger("irci.export")


# --------------------------------------------------------------------------------------
# Excel Export
# --------------------------------------------------------------------------------------
def export_to_excel(
    df_composite: pd.DataFrame,
    df_valuation: Optional[pd.DataFrame] = None,
    df_liquidity: Optional[pd.DataFrame] = None,
    df_coverage: Optional[pd.DataFrame] = None,
    df_trust: Optional[pd.DataFrame] = None,
    recommendations: Optional[List[Dict]] = None,
    company_name: str = "Analysis",
    include_charts: bool = True
) -> io.BytesIO:
    """
    Export all IRCI data to a multi-sheet Excel workbook.

    Args:
        df_composite: Main composite DataFrame
        df_valuation: Valuation dial data
        df_liquidity: Liquidity dial data
        df_coverage: Coverage dial data
        df_trust: Trust dial data
        recommendations: Playbook recommendations
        company_name: Name for the workbook
        include_charts: Whether to include charts (requires openpyxl)

    Returns:
        BytesIO buffer containing the Excel file
    """
    output = io.BytesIO()

    try:
        with pd.ExcelWriter(output, engine='openpyxl') as writer:
            # Summary sheet
            if df_composite is not None and not df_composite.empty:
                summary_cols = ['ticker', 'quarter_end', 'irci_composite_pct',
                               'valuation_pct', 'liquidity_pct', 'coverage_pct', 'sentiment_pct']
                summary_cols = [c for c in summary_cols if c in df_composite.columns]
                df_composite[summary_cols].to_excel(writer, sheet_name='IRCI Summary', index=False)

            # Individual dial sheets
            if df_valuation is not None and not df_valuation.empty:
                df_valuation.to_excel(writer, sheet_name='Valuation Details', index=False)

            if df_liquidity is not None and not df_liquidity.empty:
                df_liquidity.to_excel(writer, sheet_name='Liquidity Details', index=False)

            if df_coverage is not None and not df_coverage.empty:
                df_coverage.to_excel(writer, sheet_name='Coverage Details', index=False)

            if df_trust is not None and not df_trust.empty:
                df_trust.to_excel(writer, sheet_name='Trust Details', index=False)

            # Recommendations sheet
            if recommendations:
                rec_df = pd.DataFrame(recommendations)
                rec_df.to_excel(writer, sheet_name='Recommendations', index=False)

            # Metadata sheet
            metadata = pd.DataFrame([{
                'Generated': datetime.now().strftime('%Y-%m-%d %H:%M:%S'),
                'Company': company_name,
                'Tickers Analyzed': ', '.join(df_composite['ticker'].unique()) if df_composite is not None else 'N/A',
                'Quarters Covered': len(df_composite['quarter_end'].unique()) if df_composite is not None and 'quarter_end' in df_composite.columns else 0,
                'Source': 'IRCI - Investor Relations Composite Index'
            }])
            metadata.to_excel(writer, sheet_name='Metadata', index=False)

            # Add charts if requested and openpyxl available
            if include_charts and df_composite is not None and not df_composite.empty:
                try:
                    _add_excel_charts(writer, df_composite)
                except Exception as e:
                    log.warning(f"Could not add charts to Excel: {e}")

    except Exception as e:
        log.error(f"Excel export failed: {e}")
        raise

    output.seek(0)
    return output


def _add_excel_charts(writer, df_composite: pd.DataFrame):
    """Add charts to Excel workbook using openpyxl"""
    try:
        from openpyxl.chart import LineChart, Reference

        workbook = writer.book
        worksheet = workbook['IRCI Summary']

        # Create line chart for IRCI over time
        chart = LineChart()
        chart.title = "IRCI Score Trend"
        chart.y_axis.title = "Score (%)"
        chart.x_axis.title = "Quarter"

        # Find column indices
        cols = list(df_composite.columns)
        if 'irci_composite_pct' in cols:
            irci_col = cols.index('irci_composite_pct') + 1
            data = Reference(worksheet, min_col=irci_col, min_row=1, max_row=len(df_composite) + 1)
            chart.add_data(data, titles_from_data=True)

            # Position chart
            worksheet.add_chart(chart, "J2")

    except ImportError:
        log.warning("openpyxl charts not available")
    except Exception as e:
        log.warning(f"Chart creation failed: {e}")


# --------------------------------------------------------------------------------------
# CSV Batch Export
# --------------------------------------------------------------------------------------
def export_to_csv_batch(
    df_composite: pd.DataFrame,
    df_valuation: Optional[pd.DataFrame] = None,
    df_liquidity: Optional[pd.DataFrame] = None,
    df_coverage: Optional[pd.DataFrame] = None,
    df_trust: Optional[pd.DataFrame] = None,
    recommendations: Optional[List[Dict]] = None,
    include_individual_tickers: bool = True
) -> io.BytesIO:
    """
    Export all IRCI data to multiple CSV files in a ZIP archive.

    Args:
        df_composite: Main composite DataFrame
        df_valuation: Valuation dial data
        df_liquidity: Liquidity dial data
        df_coverage: Coverage dial data
        df_trust: Trust dial data
        recommendations: Playbook recommendations
        include_individual_tickers: Create separate CSVs per ticker

    Returns:
        BytesIO buffer containing the ZIP file
    """
    output = io.BytesIO()

    with zipfile.ZipFile(output, 'w', zipfile.ZIP_DEFLATED) as zf:
        # Main composite
        if df_composite is not None and not df_composite.empty:
            csv_buffer = io.StringIO()
            df_composite.to_csv(csv_buffer, index=False)
            zf.writestr('irci_composite.csv', csv_buffer.getvalue())

            # Individual ticker files
            if include_individual_tickers:
                for ticker in df_composite['ticker'].unique():
                    ticker_df = df_composite[df_composite['ticker'] == ticker]
                    csv_buffer = io.StringIO()
                    ticker_df.to_csv(csv_buffer, index=False)
                    zf.writestr(f'tickers/{ticker}_composite.csv', csv_buffer.getvalue())

        # Dial details
        for name, df in [
            ('valuation', df_valuation),
            ('liquidity', df_liquidity),
            ('coverage', df_coverage),
            ('trust', df_trust)
        ]:
            if df is not None and not df.empty:
                csv_buffer = io.StringIO()
                df.to_csv(csv_buffer, index=False)
                zf.writestr(f'{name}_details.csv', csv_buffer.getvalue())

        # Recommendations
        if recommendations:
            csv_buffer = io.StringIO()
            pd.DataFrame(recommendations).to_csv(csv_buffer, index=False)
            zf.writestr('recommendations.csv', csv_buffer.getvalue())

        # README
        readme = """IRCI Export Package
==================

Files included:
- irci_composite.csv: Main IRCI scores for all tickers
- valuation_details.csv: Valuation dial metrics
- liquidity_details.csv: Liquidity dial metrics
- coverage_details.csv: Coverage dial metrics
- trust_details.csv: Trust/Sentiment dial metrics
- recommendations.csv: Playbook recommendations
- tickers/: Individual ticker data files

Generated: {timestamp}
Source: IRCI - Investor Relations Composite Index

For more information, visit: https://ircibeta.streamlit.app
""".format(timestamp=datetime.now().strftime('%Y-%m-%d %H:%M:%S'))
        zf.writestr('README.txt', readme)

    output.seek(0)
    return output


# --------------------------------------------------------------------------------------
# PowerPoint Export
# --------------------------------------------------------------------------------------
def export_to_powerpoint(
    df_composite: pd.DataFrame,
    recommendations: Optional[List[Dict]] = None,
    company_name: str = "IRCI Analysis",
    include_charts: bool = True
) -> io.BytesIO:
    """
    Export IRCI analysis to PowerPoint presentation.

    Args:
        df_composite: Main composite DataFrame
        recommendations: Playbook recommendations
        company_name: Title for the presentation
        include_charts: Whether to include charts

    Returns:
        BytesIO buffer containing the PPTX file
    """
    try:
        from pptx import Presentation
        from pptx.util import Inches, Pt
        from pptx.dml.color import RgbColor
        from pptx.enum.text import PP_ALIGN
    except ImportError:
        log.error("python-pptx not installed. Run: pip install python-pptx")
        raise ImportError("python-pptx required for PowerPoint export")

    prs = Presentation()
    prs.slide_width = Inches(13.333)
    prs.slide_height = Inches(7.5)

    # Title slide
    slide_layout = prs.slide_layouts[6]  # Blank
    slide = prs.slides.add_slide(slide_layout)

    title = slide.shapes.add_textbox(Inches(0.5), Inches(2.5), Inches(12), Inches(1))
    title_frame = title.text_frame
    title_para = title_frame.paragraphs[0]
    title_para.text = company_name
    title_para.font.size = Pt(44)
    title_para.font.bold = True
    title_para.alignment = PP_ALIGN.CENTER

    subtitle = slide.shapes.add_textbox(Inches(0.5), Inches(3.8), Inches(12), Inches(0.5))
    sub_frame = subtitle.text_frame
    sub_para = sub_frame.paragraphs[0]
    sub_para.text = f"IRCI Analysis Report | {datetime.now().strftime('%B %Y')}"
    sub_para.font.size = Pt(24)
    sub_para.alignment = PP_ALIGN.CENTER

    # Executive Summary slide
    if df_composite is not None and not df_composite.empty:
        slide = prs.slides.add_slide(prs.slide_layouts[6])

        # Title
        title = slide.shapes.add_textbox(Inches(0.5), Inches(0.3), Inches(12), Inches(0.6))
        title.text_frame.paragraphs[0].text = "Executive Summary"
        title.text_frame.paragraphs[0].font.size = Pt(32)
        title.text_frame.paragraphs[0].font.bold = True

        # Calculate summary stats
        tickers = df_composite['ticker'].unique()
        latest = df_composite.groupby('ticker').last()

        avg_irci = latest['irci_composite_pct'].mean() if 'irci_composite_pct' in latest.columns else 50
        top_ticker = latest['irci_composite_pct'].idxmax() if 'irci_composite_pct' in latest.columns else tickers[0]
        top_score = latest['irci_composite_pct'].max() if 'irci_composite_pct' in latest.columns else 50

        # Summary box
        summary_text = f"""
Key Findings:
• {len(tickers)} companies analyzed
• Average IRCI Score: {avg_irci:.1f}%
• Top Performer: {top_ticker} ({top_score:.1f}%)

Dial Averages:
• Valuation: {latest['valuation_pct'].mean():.1f}% (vs sector)
• Liquidity: {latest['liquidity_pct'].mean():.1f}%
• Coverage: {latest['coverage_pct'].mean():.1f}%
• Trust: {latest['sentiment_pct'].mean():.1f}%
"""
        summary_box = slide.shapes.add_textbox(Inches(0.5), Inches(1.2), Inches(6), Inches(5))
        summary_box.text_frame.paragraphs[0].text = summary_text
        summary_box.text_frame.paragraphs[0].font.size = Pt(16)

        # Add score table
        _add_score_table(slide, latest, Inches(7), Inches(1.5))

    # Recommendations slide
    if recommendations:
        slide = prs.slides.add_slide(prs.slide_layouts[6])

        title = slide.shapes.add_textbox(Inches(0.5), Inches(0.3), Inches(12), Inches(0.6))
        title.text_frame.paragraphs[0].text = "Top Recommendations"
        title.text_frame.paragraphs[0].font.size = Pt(32)
        title.text_frame.paragraphs[0].font.bold = True

        y_pos = Inches(1.2)
        for i, rec in enumerate(recommendations[:5]):  # Top 5 recommendations
            rec_box = slide.shapes.add_textbox(Inches(0.5), y_pos, Inches(12), Inches(1))
            frame = rec_box.text_frame
            p = frame.paragraphs[0]
            p.text = f"{i+1}. {rec.get('action', 'Action')}"
            p.font.size = Pt(18)
            p.font.bold = True

            # Add description
            p2 = frame.add_paragraph()
            p2.text = rec.get('description', '')[:200]  # Truncate long descriptions
            p2.font.size = Pt(14)

            y_pos += Inches(1.1)

    output = io.BytesIO()
    prs.save(output)
    output.seek(0)
    return output


def _add_score_table(slide, latest_df: pd.DataFrame, left: float, top: float):
    """Add a simple score table to a PowerPoint slide"""
    try:
        from pptx.util import Inches, Pt
        from pptx.enum.text import PP_ALIGN

        rows = min(len(latest_df), 10) + 1  # Header + data rows
        cols = 5

        table = slide.shapes.add_table(rows, cols, left, top, Inches(5.5), Inches(0.4 * rows)).table

        # Header
        headers = ['Ticker', 'IRCI', 'Valuation', 'Liquidity', 'Coverage']
        for i, header in enumerate(headers):
            cell = table.cell(0, i)
            cell.text = header
            cell.text_frame.paragraphs[0].font.size = Pt(12)
            cell.text_frame.paragraphs[0].font.bold = True

        # Data
        for row_idx, (ticker, row) in enumerate(latest_df.head(10).iterrows()):
            table.cell(row_idx + 1, 0).text = str(ticker)
            table.cell(row_idx + 1, 1).text = f"{row.get('irci_composite_pct', 0):.0f}%"
            table.cell(row_idx + 1, 2).text = f"{row.get('valuation_pct', 0):.0f}%"
            table.cell(row_idx + 1, 3).text = f"{row.get('liquidity_pct', 0):.0f}%"
            table.cell(row_idx + 1, 4).text = f"{row.get('coverage_pct', 0):.0f}%"

    except Exception as e:
        log.warning(f"Could not add table to PowerPoint: {e}")


# --------------------------------------------------------------------------------------
# Portfolio Manager Export (Batch analysis)
# --------------------------------------------------------------------------------------
def export_portfolio_analysis(
    tickers: List[str],
    df_composite: pd.DataFrame,
    format: str = 'excel',
    include_recommendations: bool = True
) -> io.BytesIO:
    """
    Export comprehensive portfolio analysis for portfolio managers.

    Args:
        tickers: List of tickers in portfolio
        df_composite: Composite data for all tickers
        format: 'excel', 'csv', or 'pptx'
        include_recommendations: Include playbook recommendations

    Returns:
        BytesIO buffer with export file
    """
    # Filter to portfolio tickers
    portfolio_df = df_composite[df_composite['ticker'].isin([t.upper() for t in tickers])]

    recommendations = None
    if include_recommendations:
        try:
            from .playbook import generate_playbook
            recommendations = []
            for ticker in tickers:
                ticker_recs = generate_playbook(portfolio_df, ticker)
                for rec in ticker_recs:
                    rec['ticker'] = ticker
                    recommendations.append(rec)
        except Exception as e:
            log.warning(f"Could not generate recommendations: {e}")

    if format == 'excel':
        return export_to_excel(
            portfolio_df,
            recommendations=recommendations,
            company_name=f"Portfolio Analysis ({len(tickers)} companies)"
        )
    elif format == 'csv':
        return export_to_csv_batch(
            portfolio_df,
            recommendations=recommendations,
            include_individual_tickers=True
        )
    elif format == 'pptx':
        return export_to_powerpoint(
            portfolio_df,
            recommendations=recommendations,
            company_name=f"Portfolio IRCI Analysis"
        )
    else:
        raise ValueError(f"Unsupported format: {format}")


# --------------------------------------------------------------------------------------
# Quick Export Functions (for UI integration)
# --------------------------------------------------------------------------------------
def get_export_filename(company_name: str, format: str) -> str:
    """Generate standardized export filename"""
    timestamp = datetime.now().strftime('%Y%m%d_%H%M')
    safe_name = "".join(c for c in company_name if c.isalnum() or c in (' ', '-', '_')).strip()
    safe_name = safe_name.replace(' ', '_')[:50]

    extensions = {
        'excel': 'xlsx',
        'csv': 'zip',
        'pptx': 'pptx',
        'pdf': 'pdf'
    }

    return f"IRCI_{safe_name}_{timestamp}.{extensions.get(format, 'xlsx')}"


def get_mime_type(format: str) -> str:
    """Get MIME type for download"""
    mime_types = {
        'excel': 'application/vnd.openxmlformats-officedocument.spreadsheetml.sheet',
        'csv': 'application/zip',
        'pptx': 'application/vnd.openxmlformats-officedocument.presentationml.presentation',
        'pdf': 'application/pdf'
    }
    return mime_types.get(format, 'application/octet-stream')
